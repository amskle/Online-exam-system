"""轻量文档解析器集合。"""
from __future__ import annotations

import importlib.util
import logging
import re
from collections import defaultdict
from pathlib import Path

from config.settings import get_settings
from rag.models import ParsedBlock, ParsedDocument


logger = logging.getLogger("ai-tutor.rag.parsers")
settings = get_settings()


_PAGE_NOISE_RES = [
    re.compile(r"第\s*\d+\s*[／/]\s*\d+\s*页"),
    re.compile(r"^\d{1,3}\s*$", re.MULTILINE),
    re.compile(r"^\s*[-—}》\s]+\s*$", re.MULTILINE),
]
_OCR_FIXES = [
    (re.compile(r"(?:^|\n)\s*[lI]\.\s*(?=[^\d])"), "\n1. "),
    (re.compile(r"(?:^|\n)\s*[lI]\s+(\d)\s*\."), r"\n1\1."),
    (re.compile(r"(?:^|\n)\s*(\d)\s+(\d)\s*\."), r"\n\1\2."),
]
_VERTICAL_CHAR_BLOCK = re.compile(r"(?:^[^\n]{1,3}\n){5,}", re.MULTILINE)
_MARKDOWN_HEADING = re.compile(r"^(#{1,6})\s+(.+)$")
_NUMBERED_HEADING = re.compile(
    r"^(?:第[一二三四五六七八九十百千万0-9]+[章节篇卷]|[一二三四五六七八九十]+、|(?:\d+\.)+\d*)"
)


def _unstructured_enabled() -> bool:
    if not settings.use_unstructured:
        return False
    if importlib.util.find_spec("unstructured") is None:
        logger.warning("USE_UNSTRUCTURED=true 但未安装 unstructured，回退轻量解析")
        return False
    return True


def _element_blocks(elements) -> list[ParsedBlock]:
    blocks: list[ParsedBlock] = []
    for element in elements:
        text = str(getattr(element, "text", "") or "").strip()
        if not text:
            continue
        category = str(getattr(element, "category", "")).lower()
        if category == "pagebreak":
            continue
        if category in ("title", "header", "sectionheader"):
            blocks.append(ParsedBlock(text=text, kind="heading", level=1))
        elif category == "table":
            blocks.append(ParsedBlock(text=text, kind="table"))
        else:
            blocks.append(ParsedBlock(text=text))
    return blocks


def _clean_pdf_text(text: str) -> str:
    for pat, repl in _OCR_FIXES:
        text = pat.sub(repl, text)
    text = _VERTICAL_CHAR_BLOCK.sub("\n", text)
    for pat in _PAGE_NOISE_RES:
        text = pat.sub("", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _heading_level(line: str) -> int:
    line = line.strip()
    md = _MARKDOWN_HEADING.match(line)
    if md:
        return len(md.group(1))
    if _NUMBERED_HEADING.match(line) and len(line) < 80:
        dots = line.split(".")[0].count(".")
        return min(6, dots + 1)
    return 0


def _split_paragraphs(text: str) -> list[str]:
    parts = re.split(r"\n\s*\n", text)
    return [p.strip() for p in parts if p.strip()]


def _parse_text_lines(text: str) -> list[ParsedBlock]:
    blocks: list[ParsedBlock] = []
    paragraph: list[str] = []

    def flush_paragraph():
        if paragraph:
            content = "\n".join(paragraph).strip()
            if content:
                blocks.append(ParsedBlock(text=content))
            paragraph.clear()

    for raw_line in text.splitlines():
        line = raw_line.rstrip()
        level = _heading_level(line)
        if level:
            flush_paragraph()
            blocks.append(ParsedBlock(text=line.strip(), kind="heading", level=level))
        elif line.strip():
            paragraph.append(line.strip())
        else:
            flush_paragraph()
    flush_paragraph()
    return blocks


def _parse_markdown(text: str) -> list[ParsedBlock]:
    blocks: list[ParsedBlock] = []
    paragraph: list[str] = []
    code_lines: list[str] = []
    code_lang = ""
    in_code = False

    def flush_paragraph():
        if paragraph:
            content = "\n".join(paragraph).strip()
            if content:
                blocks.append(ParsedBlock(text=content))
            paragraph.clear()

    def flush_code():
        if code_lines:
            content = "\n".join(code_lines).strip()
            if content:
                blocks.append(ParsedBlock(text=content, kind="code"))
            code_lines.clear()

    for raw_line in text.splitlines():
        stripped = raw_line.strip()
        if stripped.startswith("```") or stripped.startswith("~~~"):
            if in_code:
                in_code = False
                flush_code()
                code_lang = ""
            else:
                flush_paragraph()
                in_code = True
                code_lang = stripped[3:].strip()
            continue
        if in_code:
            code_lines.append(raw_line)
            continue

        md = _MARKDOWN_HEADING.match(raw_line)
        if md:
            flush_paragraph()
            blocks.append(ParsedBlock(text=md.group(2).strip(), kind="heading", level=len(md.group(1))))
        elif stripped:
            paragraph.append(stripped)
        else:
            flush_paragraph()
    flush_paragraph()
    flush_code()
    return blocks


def _pdf_paragraphs(text: str) -> list[ParsedBlock]:
    cleaned = _clean_pdf_text(text)
    blocks: list[ParsedBlock] = []
    for para in _split_paragraphs(cleaned):
        level = _heading_level(para)
        if level and len(para) < 80:
            blocks.append(ParsedBlock(text=para, kind="heading", level=level))
        else:
            blocks.append(ParsedBlock(text=para))
    return blocks


def parse_pdf(file_path: str) -> ParsedDocument:
    if _unstructured_enabled():
        try:
            from unstructured.partition.pdf import partition_pdf

            elements = partition_pdf(file_path, strategy="auto")
            return ParsedDocument(format="pdf", blocks=_element_blocks(elements))
        except Exception as e:
            logger.warning("Unstructured PDF 解析失败，回退 PyMuPDF: %s", e)

    import fitz

    doc = fitz.open(file_path)
    toc = doc.get_toc(simple=True) or []
    toc_by_page: dict[int, list[tuple[int, str]]] = defaultdict(list)
    for level, title, page in toc:
        toc_by_page[int(page)].append((int(level), str(title).strip()))

    blocks: list[ParsedBlock] = []
    for page_index, page in enumerate(doc, start=1):
        for level, title in sorted(toc_by_page.get(page_index, []), key=lambda item: item[0]):
            blocks.append(ParsedBlock(text=title, kind="heading", level=level, page_range=str(page_index)))
        page_blocks = _pdf_paragraphs(page.get_text("text"))
        for block in page_blocks:
            block.page_range = block.page_range or str(page_index)
            blocks.append(block)

    metadata = {}
    if doc.metadata and doc.metadata.get("creationDate"):
        metadata["created_at"] = doc.metadata["creationDate"]
    doc.close()
    return ParsedDocument(format="pdf", blocks=blocks, metadata=metadata)


def parse_txt(file_path: str) -> ParsedDocument:
    raw = Path(file_path).read_bytes()
    text = None
    for encoding in ("utf-8", "gbk"):
        try:
            text = raw.decode(encoding)
            break
        except UnicodeDecodeError:
            continue
    if text is None:
        text = raw.decode("utf-8", errors="replace")
    return ParsedDocument(format="txt", blocks=_parse_text_lines(text))


def parse_markdown(file_path: str) -> ParsedDocument:
    text = Path(file_path).read_text(encoding="utf-8", errors="replace")
    blocks = _parse_markdown(text)
    return ParsedDocument(format="md", blocks=blocks)


def _docx_heading_level(style_name: str | None) -> int:
    name = (style_name or "").lower()
    match = re.search(r"heading\s*(\d)|标题\s*(\d)", name)
    if match:
        return int(next(g for g in match.groups() if g))
    if "title" in name or "标题" in name:
        return 1
    return 0


def parse_docx(file_path: str) -> ParsedDocument:
    if _unstructured_enabled():
        try:
            from unstructured.partition.docx import partition_docx

            elements = partition_docx(file_path)
            return ParsedDocument(format="docx", blocks=_element_blocks(elements))
        except Exception as e:
            logger.warning("Unstructured DOCX 解析失败，回退 python-docx: %s", e)

    from docx import Document
    from docx.table import Table
    from docx.text.paragraph import Paragraph

    doc = Document(file_path)
    blocks: list[ParsedBlock] = []

    for child in doc.element.body.iterchildren():
        tag = child.tag.rsplit("}", 1)[-1]
        if tag == "p":
            para = Paragraph(child, doc)
            text = para.text.strip()
            if not text:
                continue
            level = _docx_heading_level(para.style.name if para.style else None)
            kind = "heading" if level else "paragraph"
            blocks.append(ParsedBlock(text=text, kind=kind, level=level))
        elif tag == "tbl":
            table = Table(child, doc)
            for row in table.rows:
                for cell in row.cells:
                    for para in cell.paragraphs:
                        text = para.text.strip()
                        if text:
                            blocks.append(ParsedBlock(text=text, kind="table"))

    metadata = {}
    props = doc.core_properties
    if getattr(props, "created", None):
        metadata["created_at"] = props.created.timestamp()
    if getattr(props, "modified", None):
        metadata["modified_at"] = props.modified.timestamp()
    return ParsedDocument(format="docx", blocks=blocks, metadata=metadata)


def _shape_text(shape) -> list[str]:
    texts: list[str] = []
    if not getattr(shape, "has_text_frame", False):
        return texts
    for para in shape.text_frame.paragraphs:
        text = para.text.strip()
        if text:
            texts.append(text)
    return texts


def parse_pptx(file_path: str) -> ParsedDocument:
    if _unstructured_enabled():
        try:
            from unstructured.partition.pptx import partition_pptx

            elements = partition_pptx(file_path)
            return ParsedDocument(format="pptx", blocks=_element_blocks(elements))
        except Exception as e:
            logger.warning("Unstructured PPTX 解析失败，回退 python-pptx: %s", e)

    from pptx import Presentation

    presentation = Presentation(file_path)
    blocks: list[ParsedBlock] = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        page = str(slide_index)
        title_shape = getattr(slide, "shapes", None).title if getattr(slide, "shapes", None) else None
        title = title_shape.text.strip() if title_shape is not None and title_shape.text.strip() else ""
        if title:
            blocks.append(ParsedBlock(text=title, kind="heading", level=1, page_range=page))

        for shape in slide.shapes:
            if shape == title_shape:
                continue
            for text in _shape_text(shape):
                blocks.append(ParsedBlock(text=text, kind="paragraph", page_range=page))
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        text = cell.text.strip()
                        if text:
                            blocks.append(ParsedBlock(text=text, kind="table", page_range=page))

        if slide.has_notes_slide:
            notes_text = slide.notes_slide.notes_text_frame.text.strip()
            if notes_text:
                for part in _split_paragraphs(notes_text):
                    blocks.append(ParsedBlock(text=part, kind="note", page_range=page))

    metadata = {}
    props = presentation.core_properties
    if getattr(props, "created", None):
        metadata["created_at"] = props.created.timestamp()
    if getattr(props, "modified", None):
        metadata["modified_at"] = props.modified.timestamp()
    return ParsedDocument(format="pptx", blocks=blocks, metadata=metadata)


def parse_document(file_path: str, fmt: str) -> ParsedDocument:
    parsers = {
        "pdf": parse_pdf,
        "txt": parse_txt,
        "md": parse_markdown,
        "docx": parse_docx,
        "pptx": parse_pptx,
    }
    parser = parsers.get(fmt)
    if not parser:
        raise ValueError(f"不支持的文档格式: {fmt}")
    return parser(file_path)

"""RAG 格式识别、分类与分块策略测试。"""
import zipfile

import pytest

from rag.chunkers.classifier import classify
from rag.chunkers.strategies import (
    ChapterChunker,
    NotesChunker,
    QuestionBankChunker,
    estimate_tokens,
    split_text_into_chunks,
)
from rag.document_loader import DocumentLoader
from rag.keywords import extract_keyword_terms
from rag.loaders.format_detector import detect_format
from rag.loaders.parsers import parse_document
from rag.models import DocumentChunk, ParsedBlock
from rag.vector_store import VectorStore


def test_detect_format_text_markdown_office_and_unknown(tmp_path):
    txt = tmp_path / "a.txt"
    txt.write_text("普通文本", encoding="utf-8")
    assert detect_format(txt) == "txt"

    md = tmp_path / "b.md"
    md.write_text("# 标题\n\n正文", encoding="utf-8")
    assert detect_format(md) == "md"

    docx = tmp_path / "c.docx"
    with zipfile.ZipFile(docx, "w") as zf:
        zf.writestr("[Content_Types].xml", "<Types/>")
        zf.writestr("word/document.xml", "<w:document/>")
    assert detect_format(docx) == "docx"

    pptx = tmp_path / "d.pptx"
    with zipfile.ZipFile(pptx, "w") as zf:
        zf.writestr("[Content_Types].xml", "<Types/>")
        zf.writestr("ppt/presentation.xml", "<p:presentation/>")
    assert detect_format(pptx) == "pptx"

    unknown = tmp_path / "e.bin"
    unknown.write_bytes(b"\x00\x01\x02binary")
    with pytest.raises(ValueError):
        detect_format(unknown)


def test_classifier_distinguishes_document_types():
    question_text = "1. 1+1=?\n答案：2\n\n2. 什么是进程？\n正确答案：运行中的程序"
    assert classify([], question_text) == "question_bank"

    chapter_blocks = [
        ParsedBlock(text="第一章 概述", kind="heading", level=1),
        ParsedBlock(text="正文"),
        ParsedBlock(text="1.1 小节", kind="heading", level=2),
        ParsedBlock(text="小节正文"),
    ]
    assert classify(chapter_blocks, "第一章 概述\n正文\n1.1 小节\n小节正文") == "chapter"

    note_blocks = [ParsedBlock(text="这是一个零散待办\n另一个短段落")]
    assert classify(note_blocks, "这是一个零散待办\n另一个短段落") == "scattered_notes"


def test_question_bank_chunker_keeps_index_and_splits_long_question():
    question = "1. 很长的一道题。\n" + "知识点解析内容。" * 30 + "\n答案：B"
    meta = {"source_file": "q.txt", "subject": "CS", "format": "txt", "structure_type": "question_bank"}
    chunks = QuestionBankChunker.chunk(question, meta, max_tokens=80, max_chars=120)
    assert len(chunks) >= 2
    assert {c.metadata["question_index"] for c in chunks} == {1}
    assert chunks[0].metadata["chunk_type"] == "question"
    assert chunks[1].metadata["chunk_type"] == "question_part"


def test_chapter_chunker_merges_short_sections_within_parent():
    blocks = [
        ParsedBlock(text="第一章", kind="heading", level=1),
        ParsedBlock(text="章引言"),
        ParsedBlock(text="1.1", kind="heading", level=2),
        ParsedBlock(text="小节一"),
        ParsedBlock(text="1.2", kind="heading", level=2),
        ParsedBlock(text="小节二"),
        ParsedBlock(text="第二章", kind="heading", level=1),
        ParsedBlock(text="另一章"),
    ]
    meta = {"source_file": "book.md", "subject": "OS", "format": "md", "structure_type": "chapter"}
    chunks = ChapterChunker.chunk(
        blocks,
        meta,
        max_tokens=400,
        max_chars=500,
        overlap_ratio=0.15,
        min_tokens=50,
    )
    titles = {c.metadata["section_title"] for c in chunks}
    assert "1.1等" in titles
    assert "第二章" in titles
    for chunk in chunks:
        if chunk.metadata["section_title"] == "1.1等":
            assert chunk.metadata["section_path"] == "第一章 / 1.1"


def test_notes_chunker_overlap_and_short_chunk_merge():
    paragraph = "人工智能算法说明，用于测试重叠窗口。"
    text = "\n\n".join(f"第{i}段：{paragraph}" for i in range(12))
    blocks = [ParsedBlock(text=text)]
    meta = {"source_file": "notes.md", "subject": "AI", "format": "md", "structure_type": "scattered_notes"}
    chunks = NotesChunker.chunk(
        blocks,
        meta,
        max_tokens=80,
        max_chars=500,
        overlap_ratio=0.25,
        soft_tokens=50,
        min_tokens=10,
    )
    assert len(chunks) > 1
    assert "第1段" in chunks[1].content

    short_chunks = [
        DocumentChunk(content="短内容甲", metadata={"section_path": "A"}),
        DocumentChunk(content="短内容乙", metadata={"section_path": "A"}),
        DocumentChunk(content="短内容丙", metadata={"section_path": "B"}),
    ]
    merged = NotesChunker._merge_short_chunks(short_chunks, min_tokens=50)
    assert len(merged) == 2
    assert "短内容甲" in merged[0].content and "短内容乙" in merged[0].content
    assert merged[1].content == "短内容丙"


def test_document_loader_markdown_pipeline(tmp_path):
    md = tmp_path / "textbook.md"
    md.write_text(
        "# 第一章 计算机网络\n\n网络是计算机通信的基础。\n\n## 1.1 分层模型\n\nOSI 模型分为七层。\n",
        encoding="utf-8",
    )
    result = DocumentLoader.load_and_chunk_detail(
        str(md),
        subject_name="网络",
        modified_at=1234567000,
        uploaded_at=1234568000,
    )
    assert result.format == "md"
    assert result.structure_type == "chapter"
    assert result.chunking_strategy == "chapter"
    assert result.chunks[0].metadata["source_file"] == "textbook.md"
    assert result.chunks[0].metadata["section_path"]
    assert result.chunks[0].metadata["modified_at"] == 1234567000


def test_token_estimate_and_chunk_hard_cap():
    assert estimate_tokens("人工智能") >= 4
    parts = split_text_into_chunks(
        "句子一。" + "很长很长的中文内容。" * 20,
        max_tokens=60,
        max_chars=80,
        overlap_ratio=0.0,
    )
    assert len(parts) > 1
    assert all(len(part) <= 80 for part in parts)


def test_keyword_terms_preserve_case():
    terms = extract_keyword_terms("Python 机器学习")
    assert any(term in ("Python", "python") for term in terms)
    assert "机器" in terms
    assert "学习" in terms


def test_docx_and_pptx_parsers_keep_heading_structure(tmp_path):
    from docx import Document as DocxDocument

    docx_path = tmp_path / "book.docx"
    doc = DocxDocument()
    doc.add_heading("第一章 概述", level=1)
    doc.add_paragraph("这是正文段落")
    doc.save(docx_path)
    parsed_docx = parse_document(str(docx_path), "docx")
    assert any(b.kind == "heading" and b.level == 1 for b in parsed_docx.blocks)

    from pptx import Presentation

    pptx_path = tmp_path / "slides.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "第一章 概述"
    slide.placeholders[1].text = "这是正文段落"
    prs.save(pptx_path)
    parsed_pptx = parse_document(str(pptx_path), "pptx")
    assert any(b.kind == "heading" and b.level == 1 for b in parsed_pptx.blocks)


def test_vector_store_search_does_not_require_unbound_self():
    class FakeCollection:
        def query(self, **kwargs):
            return {
                "ids": [["a"]],
                "documents": [["Thymeleaf 是模板引擎"]],
                "metadatas": [[{"source_file": "3. Thymeleaf.pdf"}]],
                "distances": [[0.1]],
            }

    store = object.__new__(VectorStore)
    results = store._search(FakeCollection(), [0.1, 0.2, 0.3], top_k=1)
    assert results[0]["id"] == "a"
    assert results[0]["metadata"]["source_file"] == "3. Thymeleaf.pdf"


def test_document_loader_keeps_original_source_name(tmp_path):
    md = tmp_path / "temp_upload.md"
    md.write_text("# Thymeleaf\n\nThymeleaf 是 Java 模板引擎。", encoding="utf-8")
    result = DocumentLoader.load_and_chunk_detail(
        str(md),
        subject_name="Java",
        source_name="3. Thymeleaf.pdf",
    )
    assert result.chunks[0].metadata["source_file"] == "3. Thymeleaf.pdf"
